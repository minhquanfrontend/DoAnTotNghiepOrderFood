#!/usr/bin/env python3
"""
Script to create PowerPoint presentation for Food Delivery System project
Run: pip install python-pptx
Then: python create_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    # Create presentation object
    prs = Presentation()
    
    # Define colors
    primary_color = RGBColor(52, 152, 219)  # Blue
    secondary_color = RGBColor(44, 62, 80)  # Dark blue
    accent_color = RGBColor(231, 76, 60)    # Red
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "HỆ THỐNG GIAO HÀNG THỰC PHẨM"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    subtitle.text = """Ứng dụng di động và hệ thống quản lý toàn diện
Tích hợp AI và công nghệ hiện đại

GVHD: [Tên Giảng viên Hướng dẫn]
SVTH: [Tên Sinh viên Thực hiện]
Lớp: [Tên Lớp] - MSSV: [Mã số sinh viên]"""
    
    # Slide 2: Problem Statement
    slide_layout = prs.slide_layouts[1]  # Title and content
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "ĐẶT VẤN ĐỀ"
    title.text_frame.paragraphs[0].font.color.rgb = secondary_color
    
    content.text = """• Vấn đề hiện tại: Khó khăn trong việc đặt món ăn trực tuyến với trải nghiệm người dùng kém
• Thiếu hệ thống: Quản lý đơn hàng, theo dõi giao hàng thời gian thực
• Không có AI: Thiếu tính năng gợi ý món ăn thông minh
• Quản lý phức tạp: Nhà hàng khó quản lý menu và đơn hàng
• Thanh toán: Thiếu tích hợp thanh toán đa dạng và an toàn"""
    
    # Slide 3: Objectives
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "MỤC TIÊU DỰ ÁN"
    title.text_frame.paragraphs[0].font.color.rgb = secondary_color
    
    content.text = """• Xây dựng ứng dụng di động: Giao diện thân thiện cho khách hàng, nhà hàng và shipper
• Hệ thống backend mạnh mẽ: API RESTful với Django, quản lý dữ liệu hiệu quả
• Tích hợp AI: Chatbot thông minh, gợi ý món ăn cá nhân hóa
• Theo dõi thời gian thực: GPS tracking, cập nhật trạng thái đơn hàng
• Thanh toán đa dạng: Tích hợp Stripe, ví điện tử
• Dashboard quản trị: Thống kê, báo cáo chi tiết"""
    
    # Slide 4: Technology Stack
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "CÔNG NGHỆ SỬ DỤNG"
    title.text_frame.paragraphs[0].font.color.rgb = secondary_color
    
    content.text = """Frontend Mobile:
• React Native, Expo, React Navigation

Backend API:
• Django, Django REST Framework, JWT Authentication

Database:
• SQLite (Development), MySQL (Production)

AI & Machine Learning:
• OpenAI GPT, Scikit-learn, Pandas, NumPy

Maps & Location:
• React Native Maps, Expo Location, Geopy

Payment & Others:
• Stripe, Celery, Redis, Django CORS"""
    
    # Slide 5: System Architecture
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "KIẾN TRÚC HỆ THỐNG"
    title.text_frame.paragraphs[0].font.color.rgb = secondary_color
    
    # Add architecture diagram as text (you can replace with actual diagram)
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(5)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = """KIẾN TRÚC 3 TẦNG:

1. PRESENTATION LAYER
   • Mobile App (React Native)
   • Admin Dashboard

2. BUSINESS LOGIC LAYER
   • Django REST API
   • Authentication Service
   • AI Features (OpenAI)
   • Order Management
   • Payment Processing

3. DATA LAYER
   • MySQL Database
   • Redis Cache
   • File Storage"""
    
    # Slide 6: Use Cases
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "SƠ ĐỒ USE CASE"
    title.text_frame.paragraphs[0].font.color.rgb = secondary_color
    
    content.text = """KHÁCH HÀNG:
• Đăng ký/Đăng nhập
• Tìm kiếm món ăn
• Đặt hàng, Thanh toán
• Theo dõi đơn hàng
• Chat với AI

NHÀ HÀNG:
• Quản lý menu
• Nhận đơn hàng
• Cập nhật trạng thái
• Xem thống kê

SHIPPER:
• Nhận đơn giao
• Định vị GPS
• Cập nhật vị trí
• Hoàn thành giao hàng"""
    
    # Slide 7: AI Algorithm
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "THUẬT TOÁN AI"
    title.text_frame.paragraphs[0].font.color.rgb = secondary_color
    
    content.text = """HỆ THỐNG GỢI Ý MÓN ĂN:
• TF-IDF Vectorization: Phân tích mô tả món ăn
• Cosine Similarity: Tính độ tương đồng giữa các món
• Collaborative Filtering: Dựa trên lịch sử đặt hàng

CHATBOT AI:
• OpenAI GPT Integration: Xử lý ngôn ngữ tự nhiên
• Context Awareness: Hiểu ngữ cảnh cuộc hội thoại
• Food Recommendation: Gợi ý dựa trên sở thích

THUẬT TOÁN TỐI ƯU:
• Route Optimization: Tối ưu đường đi cho shipper
• Load Balancing: Phân bổ đơn hàng hiệu quả"""
    
    # Slide 8: Database Design
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "THIẾT KẾ CSDL"
    title.text_frame.paragraphs[0].font.color.rgb = secondary_color
    
    content.text = """USERS & AUTHENTICATION:
• User, UserProfile, EmailVerification

RESTAURANT MANAGEMENT:
• Restaurant, Food, Category, RestaurantLocation

ORDER SYSTEM:
• Order, OrderItem, OrderStatus, OrderTracking

AI FEATURES:
• UserPreference, FoodRecommendation, ChatSession, ChatMessage

PAYMENT & WALLET:
• Payment, Transaction, Wallet, WalletTransaction

NOTIFICATIONS:
• Notification, NotificationTemplate, UserNotification"""
    
    # Slide 9: Key Features
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "TÍNH NĂNG CHÍNH"
    title.text_frame.paragraphs[0].font.color.rgb = secondary_color
    
    content.text = """🤖 AI CHATBOT:
• Tư vấn món ăn thông minh
• Xử lý đơn hàng bằng ngôn ngữ tự nhiên

📍 REAL-TIME TRACKING:
• Theo dõi shipper và đơn hàng trên bản đồ thời gian thực

💳 MULTI PAYMENT:
• Stripe, ví điện tử, thanh toán khi nhận hàng

🎯 SMART RECOMMENDATION:
• Gợi ý món ăn dựa trên AI và machine learning

📊 ANALYTICS DASHBOARD:
• Thống kê doanh thu, đơn hàng, hiệu suất"""
    
    # Slide 10: Implementation Results
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "KẾT QUẢ THỰC HIỆN"
    title.text_frame.paragraphs[0].font.color.rgb = secondary_color
    
    content.text = """THỐNG KÊ DỰ ÁN:
• 15+ Modules
• 50+ API Endpoints  
• 30+ Screens
• 3 User Roles

TÍNH NĂNG ĐÃ HOÀN THÀNH:
✅ Hệ thống đăng ký/đăng nhập với xác thực email
✅ Quản lý menu và đơn hàng cho nhà hàng
✅ Ứng dụng mobile đa nền tảng (iOS/Android)
✅ Tích hợp AI chatbot và gợi ý món ăn
✅ Hệ thống thanh toán Stripe
✅ Theo dõi GPS và bản đồ thời gian thực
✅ Dashboard quản trị với thống kê"""
    
    # Slide 11: Performance
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "HIỆU SUẤT & GIAO DIỆN"
    title.text_frame.paragraphs[0].font.color.rgb = secondary_color
    
    content.text = """HIỆU SUẤT HỆ THỐNG:
• Response Time: API < 200ms, Mobile App < 1s load time
• Scalability: Hỗ trợ 1000+ users đồng thời
• Background Tasks: Celery cho xử lý bất đồng bộ

GIAO DIỆN ỨNG DỤNG:
• Material Design: Giao diện hiện đại, thân thiện
• Responsive: Tương thích đa thiết bị
• Dark/Light Mode: Hỗ trợ chế độ sáng/tối
• Accessibility: Tuân thủ tiêu chuẩn truy cập

CÔNG NGHỆ UI/UX:
• React Native Paper cho components
• Vector Icons cho biểu tượng
• Charts và graphs cho thống kê"""
    
    # Slide 12: Conclusion
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "KẾT LUẬN & HƯỚNG PHÁT TRIỂN"
    title.text_frame.paragraphs[0].font.color.rgb = secondary_color
    
    content.text = """KẾT LUẬN:
✅ Xây dựng thành công hệ thống giao hàng thực phẩm hoàn chỉnh
✅ Tích hợp AI và machine learning hiệu quả
✅ Ứng dụng mobile đa nền tảng với UX/UI tốt
✅ Hệ thống backend mạnh mẽ, scalable

HƯỚNG PHÁT TRIỂN:
🚀 Triển khai lên cloud (AWS/Google Cloud)
🤖 Nâng cấp AI với deep learning
📱 Phát triển web app cho admin
🔄 Tích hợp thêm phương thức thanh toán
📊 Business Intelligence và Analytics nâng cao
🌐 Mở rộng ra thị trường quốc tế

CẢM ƠN QUÝ THẦY CÔ ĐÃ LẮNG NGHE!"""
    
    # Save presentation
    prs.save('Food_Delivery_System_Presentation.pptx')
    print("✅ Đã tạo thành công file Food_Delivery_System_Presentation.pptx")
    print("📁 File được lưu tại thư mục hiện tại")
    print("🎯 Bạn có thể mở file bằng Microsoft PowerPoint hoặc Google Slides")

if __name__ == "__main__":
    try:
        create_presentation()
    except ImportError:
        print("❌ Lỗi: Chưa cài đặt thư viện python-pptx")
        print("📦 Chạy lệnh: pip install python-pptx")
        print("🔄 Sau đó chạy lại script này")
    except Exception as e:
        print(f"❌ Lỗi: {str(e)}")
